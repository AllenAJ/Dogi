#!/usr/bin/env python3
"""Build Dogi hackathon pitch deck (.pptx)."""

from __future__ import annotations

import os
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "slides" / "Dogi-Hackathon-Deck.pptx"
MASCOT = ROOT / "public" / "mascot.gif"

BG = RGBColor(0xFD, 0xF7, 0xF0)
SURFACE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0xFF, 0xDD, 0x00)
DARK = RGBColor(0x0D, 0x0C, 0x22)
MUTED = RGBColor(0x6F, 0x6A, 0x5E)
BORDER = RGBColor(0xED, 0xE3, 0xD5)


def set_slide_bg(slide, color: RGBColor = BG) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_dots(slide, left=0, top=0, width=13.333, height=7.5, cols=28, rows=16) -> None:
    dx = width / cols
    dy = height / rows
    for r in range(rows):
        for c in range(cols):
            shape = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(left + c * dx + dx * 0.35),
                Inches(top + r * dy + dy * 0.35),
                Inches(0.045),
                Inches(0.045),
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = BORDER
            shape.line.fill.background()
            shape.shadow.inherit = False


def add_accent_blob(slide, left, top, size, alpha=0.22) -> None:
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(left), Inches(top), Inches(size), Inches(size)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.fill.transparency = 1 - alpha
    shape.line.fill.background()


def textbox(
    slide,
    left,
    top,
    width,
    height,
    lines: list[tuple[str, int, bool, RGBColor]],
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = valign
    for i, (text, size, bold, color) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.alignment = align
        p.space_after = Pt(8)
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = "Aptos Display" if bold and size >= 28 else "Aptos"
    return box


def add_card(slide, left, top, width, height, title: str, body: str) -> None:
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    card.fill.solid()
    card.fill.fore_color.rgb = SURFACE
    card.line.color.rgb = BORDER
    card.line.width = Pt(1)
    card.adjustments[0] = 0.12
    textbox(
        slide,
        left + 0.2,
        top + 0.18,
        width - 0.4,
        height - 0.3,
        [
            (title, 18, True, DARK),
            (body, 13, False, MUTED),
        ],
    )


def add_slide_base(slide, kicker: str | None = None) -> None:
    set_slide_bg(slide)
    add_dots(slide)
    add_accent_blob(slide, 10.5, -0.8, 3.2, 0.18)
    add_accent_blob(slide, -0.9, 5.2, 2.4, 0.14)
    if kicker:
        textbox(slide, 0.7, 0.45, 5, 0.35, [(kicker, 11, True, MUTED)])


def add_title_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_base(slide)
    if MASCOT.exists():
        slide.shapes.add_picture(str(MASCOT), Inches(5.55), Inches(0.55), height=Inches(1.15))
    textbox(
        slide,
        0.85,
        1.9,
        11.5,
        1.2,
        [("Dogi", 54, True, DARK)],
    )
    pill = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.85), Inches(3.05), Inches(4.8), Inches(0.72)
    )
    pill.fill.solid()
    pill.fill.fore_color.rgb = ACCENT
    pill.line.fill.background()
    pill.adjustments[0] = 0.5
    textbox(
        slide,
        1.0,
        3.12,
        4.5,
        0.55,
        [("Get paid from any chain", 22, True, DARK)],
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )
    textbox(
        slide,
        0.85,
        4.05,
        9.5,
        1.0,
        [
            (
                "Treat pages for creators. Payment links for everyone.",
                20,
                False,
                MUTED,
            ),
            ("UXmaxx Hackathon · Universal Accounts Track", 14, False, MUTED),
        ],
    )
    notes = slide.notes_slide.notes_text_frame
    notes.text = (
        "Open with the one-liner: Dogi is Buy Me a Coffee meets cross-chain payments. "
        "Fans and clients pay with email only. Creators always receive USDC on Arbitrum."
    )


def add_bullet_slide(
    prs: Presentation,
    kicker: str,
    title: str,
    bullets: list[str],
    notes: str,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_base(slide, kicker)
    textbox(slide, 0.85, 0.95, 11.5, 0.9, [(title, 34, True, DARK)])
    y = 2.0
    for bullet in bullets:
        row = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.85), Inches(y), Inches(11.6), Inches(0.72)
        )
        row.fill.solid()
        row.fill.fore_color.rgb = SURFACE
        row.line.color.rgb = BORDER
        row.adjustments[0] = 0.35
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(1.05), Inches(y + 0.26), Inches(0.14), Inches(0.14)
        )
        dot.fill.solid()
        dot.fill.fore_color.rgb = ACCENT
        dot.line.fill.background()
        textbox(slide, 1.35, y + 0.12, 10.8, 0.5, [(bullet, 16, False, DARK)])
        y += 0.88
    slide.notes_slide.notes_text_frame.text = notes


def add_money_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_base(slide, "BUSINESS MODEL")
    textbox(slide, 0.85, 0.95, 11.5, 0.9, [("How Dogi makes money", 34, True, DARK)])
    cards = [
        (
            "1. Routing margin",
            "A 0.5% fee on top of network costs, baked into the route quote. "
            "That is a tenth of what Buy Me a Coffee charges, and it grows with every "
            "treat and invoice paid.",
        ),
        (
            "2. Dogi Pro",
            "A monthly plan for serious creators: recurring memberships, supporter "
            "analytics, custom domains and branding. The basic page stays free forever.",
        ),
        (
            "3. Business invoicing",
            "White-label payment links and an invoicing API for shops and agencies "
            "that want to accept any token from any chain without touching a wallet.",
        ),
    ]
    for i, (title, body) in enumerate(cards):
        add_card(slide, 0.85 + i * 4.05, 2.0, 3.75, 3.6, title, body)
    badge = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.4), Inches(6.0), Inches(6.5), Inches(0.6)
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = ACCENT
    badge.line.fill.background()
    badge.adjustments[0] = 0.5
    textbox(
        slide, 3.55, 6.07, 6.2, 0.45,
        [("Free to start. We only earn when creators get paid.", 15, True, DARK)],
        align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE,
    )
    slide.notes_slide.notes_text_frame.text = (
        "Today the fee is 0% while we grow. The routing margin is a single config value, "
        "so flipping it on later needs no migration. Lead with the alignment story: "
        "unlike ad platforms or paywalls, we only make money when creators do."
    )


def add_comparison_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_base(slide, "WHY DOGI WINS")
    textbox(slide, 0.85, 0.95, 11.5, 0.9, [("Same simplicity. None of the walls.", 34, True, DARK)])

    rows = [
        ("Platform fee", "5% of every payment + processor fees", "0%. You keep what fans send"),
        ("Crypto", "Dropped crypto support entirely", "Crypto-native: any token, any chain"),
        ("Who can join", "Only where Stripe / PayPal operate", "Anyone with an email address"),
        ("Your account", "Can be banned, frozen, or held", "Self-custodial. No one can take it"),
        ("Your money", "Payout delays, chargebacks", "USDC on Arbitrum, settled on-chain"),
    ]
    col_label, col_bmc, col_dogi = 0.85, 3.35, 8.0
    header_y = 2.0
    for x, text, color in (
        (col_bmc, "Buy Me a Coffee", MUTED),
        (col_dogi, "Dogi", DARK),
    ):
        pill = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(header_y), Inches(4.4), Inches(0.5)
        )
        pill.fill.solid()
        pill.fill.fore_color.rgb = ACCENT if text == "Dogi" else SURFACE
        pill.line.color.rgb = BORDER
        pill.adjustments[0] = 0.5
        textbox(
            slide, x, header_y + 0.05, 4.4, 0.4,
            [(text, 15, True, color)],
            align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE,
        )
    y = header_y + 0.75
    for label, bmc, dogi in rows:
        textbox(slide, col_label, y + 0.1, 2.3, 0.5, [(label, 12, True, MUTED)])
        for x, text, is_dogi in ((col_bmc, bmc, False), (col_dogi, dogi, True)):
            cell = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(4.4), Inches(0.62)
            )
            cell.fill.solid()
            cell.fill.fore_color.rgb = SURFACE
            cell.line.color.rgb = BORDER
            cell.adjustments[0] = 0.3
            textbox(
                slide, x + 0.2, y + 0.08, 4.05, 0.5,
                [(text, 12, is_dogi, DARK if is_dogi else MUTED)],
                valign=MSO_ANCHOR.MIDDLE,
            )
        y += 0.78
    slide.notes_slide.notes_text_frame.text = (
        "The pitch in one table: everything people like about Buy Me a Coffee, with the platform risk removed. "
        "Sources: BMC charges a 5% platform fee; it runs on Stripe/PayPal so unsupported countries are locked out; "
        "its ToS lets it suspend accounts and hold payouts."
    )


def add_route_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_base(slide, "THE MAGIC MOMENT")
    textbox(slide, 0.85, 0.95, 11.5, 0.9, [("We make chain abstraction visible", 34, True, DARK)])
    textbox(
        slide, 0.85, 2.0, 5.6, 3.5,
        [
            ("Before every payment, the payer sees the exact", 15, False, MUTED),
            ("route the Universal Account quoted: which tokens,", 15, False, MUTED),
            ("from which chains, and what it costs. One signature.", 15, False, MUTED),
            ("", 8, False, MUTED),
            ("Judges don't have to trust that EIP-7702 chain", 15, False, MUTED),
            ("abstraction is working. They watch it work.", 15, False, MUTED),
        ],
    )
    # Recreate the app's RouteCard as shapes.
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.0), Inches(1.95), Inches(5.3), Inches(4.6)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = SURFACE
    card.line.color.rgb = BORDER
    card.line.width = Pt(1.25)
    card.adjustments[0] = 0.08
    textbox(
        slide, 7.3, 2.2, 4.7, 4.2,
        [
            ("WHERE IT COMES FROM", 11, True, MUTED),
            ("0.004 ETH on Base                    $12.40", 14, False, DARK),
            ("2.10 USDT on BNB Chain            $2.10", 14, False, DARK),
            ("0.004 SOL on Solana                  $0.62", 14, False, DARK),
            ("", 6, False, MUTED),
            ("↓  $15.00 USDC on Arbitrum", 15, True, DARK),
            ("", 6, False, MUTED),
            ("Bridging, swaps & gas                 $0.12", 12, False, MUTED),
            ("Total from your balance             $15.12", 13, True, DARK),
        ],
    )
    badge = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.15), Inches(6.05), Inches(3.0), Inches(0.42)
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = ACCENT
    badge.line.fill.background()
    badge.adjustments[0] = 0.5
    textbox(
        slide, 9.2, 6.1, 2.9, 0.35,
        [("3 CHAINS, ONE TAP", 11, True, DARK)],
        align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE,
    )
    slide.notes_slide.notes_text_frame.text = (
        "This screen is the demo's money shot: funds drawn from three chains, one signature, "
        "settled as USDC on Arbitrum. It maps directly to the 30% 'prominent use of UA + 7702' criterion."
    )


def add_two_product_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_base(slide, "TWO PRODUCTS")
    textbox(slide, 0.85, 0.95, 11.5, 0.9, [("One account. Two ways to get paid.", 34, True, DARK)])
    add_card(
        slide,
        0.85,
        2.0,
        5.55,
        3.35,
        "🦴 Treat page",
        "Buy-me-a-coffee style page. Fans pick 1, 3, or 5 treats, leave a note, and pay in one tap from any chain.",
    )
    add_card(
        slide,
        6.75,
        2.0,
        5.55,
        3.35,
        "🔗 Payment links",
        "Set an amount and note, share the link. Clients log in with email and pay with whatever they hold.",
    )
    badge = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.2), Inches(5.65), Inches(6.8), Inches(0.62)
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = ACCENT
    badge.line.fill.background()
    badge.adjustments[0] = 0.5
    textbox(
        slide,
        3.35,
        5.72,
        6.5,
        0.45,
        [("Both settle as USDC on Arbitrum", 16, True, DARK)],
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )
    slide.notes_slide.notes_text_frame.text = (
        "Stress that these are not two separate apps. Same balance, same login, same settlement rail."
    )


def add_flow_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_base(slide, "HOW IT WORKS")
    textbox(slide, 0.85, 0.95, 11.5, 0.9, [("Invisible cross-chain routing", 34, True, DARK)])
    steps = [
        ("1", "Log in with email", "Magic creates a non-custodial wallet. No app, no seed phrase."),
        ("2", "Pay with any token", "ETH on Base, SOL on Solana, USDT on BNB Chain."),
        ("3", "UA routes liquidity", "Particle Universal Account in EIP-7702 mode."),
        ("4", "Receiver gets USDC", "Settlement always lands on Arbitrum."),
    ]
    x_positions = [0.85, 3.55, 6.25, 8.95]
    for (num, title, body), x in zip(steps, x_positions):
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(2.15), Inches(2.45), Inches(3.55)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = SURFACE
        card.line.color.rgb = BORDER
        card.adjustments[0] = 0.1
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(x + 0.85), Inches(2.45), Inches(0.55), Inches(0.55)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = ACCENT
        circle.line.fill.background()
        textbox(
            slide,
            x + 0.85,
            2.52,
            0.55,
            0.4,
            [(num, 18, True, DARK)],
            align=PP_ALIGN.CENTER,
            valign=MSO_ANCHOR.MIDDLE,
        )
        textbox(
            slide,
            x + 0.15,
            3.2,
            2.15,
            2.2,
            [(title, 15, True, DARK), (body, 11, False, MUTED)],
            align=PP_ALIGN.CENTER,
        )
    slide.notes_slide.notes_text_frame.text = (
        "Walk through the payer experience first, then mention the tech only if judges ask."
    )


def add_tech_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_base(slide, "TECH STACK")
    textbox(slide, 0.85, 0.95, 11.5, 0.9, [("Built for the bounty tracks", 34, True, DARK)])
    stacks = [
        ("Magic", "Email OTP login\nEmbedded EOA wallet\nEIP-7702 signing", "Magic Labs bonus"),
        ("Particle UA", "EIP-7702 mode\nOne balance, 6 chains\nCross-chain transfer API", "Universal Accounts Track"),
        ("Arbitrum", "USDC settlement layer\nInvisible to payers\nEvery payment lands here", "Arbitrum bounty"),
    ]
    for i, (name, body, tag) in enumerate(stacks):
        x = 0.85 + i * 4.05
        add_card(slide, x, 2.0, 3.75, 3.8, name, body)
        tag_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x + 0.2), Inches(5.45), Inches(3.35), Inches(0.42)
        )
        tag_shape.fill.solid()
        tag_shape.fill.fore_color.rgb = ACCENT
        tag_shape.line.fill.background()
        tag_shape.adjustments[0] = 0.5
        textbox(
            slide,
            x + 0.25,
            5.5,
            3.25,
            0.35,
            [(tag, 11, True, DARK)],
            align=PP_ALIGN.CENTER,
            valign=MSO_ANCHOR.MIDDLE,
        )
    slide.notes_slide.notes_text_frame.text = (
        "Map each integration directly to a hackathon prize. Emphasize EIP-7702 in-place upgrade, not a new smart account."
    )


def add_demo_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_base(slide, "LIVE DEMO")
    textbox(slide, 0.85, 0.95, 11.5, 0.9, [("Try it in 2 minutes", 34, True, DARK)])
    steps = [
        "Log in with email on dogi-liard.vercel.app",
        "Create a treat page or payment link on the dashboard",
        "Open the link in incognito with a second email",
        "Review the cross-chain route, confirm in one tap",
        "Payment appears in the creator's on-chain feed",
    ]
    y = 2.0
    for i, step in enumerate(steps, start=1):
        textbox(
            slide,
            0.95,
            y,
            7.5,
            0.55,
            [(f"{i}. {step}", 15, False, DARK)],
        )
        y += 0.62
    panel = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.35), Inches(2.0), Inches(4.1), Inches(3.9)
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = SURFACE
    panel.line.color.rgb = BORDER
    panel.adjustments[0] = 0.08
    textbox(
        slide,
        8.55,
        2.2,
        3.7,
        3.5,
        [
            ("Invoice #42", 14, True, MUTED),
            ("Paid ✓", 12, True, DARK),
            ("$25.00", 36, True, DARK),
            ("settled as USDC on Arbitrum", 12, False, MUTED),
            ("", 8, False, MUTED),
            ("github.com/AllenAJ/Dogi", 13, True, DARK),
        ],
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )
    slide.notes_slide.notes_text_frame.text = (
        "Run the live demo if possible. Point judges to the GitHub repo for code review."
    )


def add_closing_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_base(slide)
    if MASCOT.exists():
        slide.shapes.add_picture(str(MASCOT), Inches(6.05), Inches(1.35), height=Inches(1.4))
    textbox(
        slide,
        0.85,
        2.2,
        11.5,
        1.0,
        [("Cross-chain payments,", 40, True, DARK), ("as easy as email.", 40, True, DARK)],
    )
    textbox(
        slide,
        0.85,
        4.0,
        10,
        1.2,
        [
            ("Dogi · github.com/AllenAJ/Dogi", 18, False, MUTED),
            ("Built with Magic · Particle Universal Accounts · Arbitrum", 14, False, MUTED),
        ],
    )
    slide.notes_slide.notes_text_frame.text = "Close with the problem you solve: normal people should not need to understand chains to get paid."


def build() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_title_slide(prs)
    add_bullet_slide(
        prs,
        "THE STATUS QUO",
        "Creator payments live in walled gardens",
        [
            "Buy Me a Coffee takes 5% of every payment, plus Stripe and PayPal processor fees",
            "It dropped crypto support entirely; fans can only pay with cards",
            "Only works where Stripe and PayPal operate, so whole regions are locked out",
            "Centralized moderation: accounts banned, payouts frozen, income gone overnight",
            "You never hold your money. Payout delays, chargebacks, platform risk",
        ],
        "Anchor on the real fear: your income living inside someone else's terms of service. "
        "BMC is the named example; the same applies to Patreon and Ko-fi.",
    )
    add_bullet_slide(
        prs,
        "THE CRYPTO GAP",
        "Crypto should fix this. Its UX says no.",
        [
            "Ask a fan to tip $5 and you're asking them to install a wallet and back up a seed phrase",
            "Which chain? Which token? One wrong answer means lost funds, so most people never try",
            "Your fan holds SOL, you want USDC on Arbitrum: bridges, swaps, and gas on N chains",
            "So creators stay in the walled gardens, and crypto tip jars stay empty",
        ],
        "The gap isn't demand, it's UX. This sets up chain abstraction as the unlock.",
    )
    add_bullet_slide(
        prs,
        "SOLUTION",
        "Dogi: get paid from any chain, with just a link",
        [
            "Payers log in with email and pay with whatever they hold, on any chain",
            "You always receive USDC on Arbitrum, in your own account. No one can freeze it",
            "0% platform fee, no sign-up gatekeeping, no country list",
            "Fully client-side: payment data lives in the URL. There is no backend to trust or subpoena",
        ],
        "This is the 30-second pitch: everything people like about BMC, minus the platform risk, "
        "powered by chain abstraction.",
    )
    add_two_product_slide(prs)
    add_flow_slide(prs)
    add_route_slide(prs)
    add_tech_slide(prs)
    add_comparison_slide(prs)
    add_money_slide(prs)
    add_demo_slide(prs)
    add_closing_slide(prs)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    return OUT


if __name__ == "__main__":
    path = build()
    print(f"Wrote {path}")
